"""Report exports — CSV downloads for MIS."""

import csv
import io
from sqlalchemy.orm import Session

from app.services import (
    inventory_intelligence,
    pendency_monitor,
    performance_analytics,
    revenue_analytics,
    tat_analysis,
)


def _to_csv(rows: list[dict]) -> str:
    if not rows:
        return ""
    fieldnames: list[str] = []
    seen: set[str] = set()
    for row in rows:
        for key in row:
            if key not in seen:
                seen.add(key)
                fieldnames.append(key)
    output = io.StringIO()
    writer = csv.DictWriter(output, fieldnames=fieldnames, extrasaction="ignore")
    writer.writeheader()
    writer.writerows(rows)
    return output.getvalue()


def export_revenue_report(db: Session, vehicle_type: str | None = None) -> str:
    states = revenue_analytics.get_revenue_by_state(db, vehicle_type)
    oems = revenue_analytics.get_revenue_by_oem(db, vehicle_type)
    portals = revenue_analytics.get_revenue_by_portal(db)
    combined = (
        [{"section": "state", **s} for s in states]
        + [{"section": "oem", **o} for o in oems]
        + [{"section": "portal", **p} for p in portals]
    )
    return _to_csv(combined)


def export_pendency_report(db: Session) -> str:
    stages = pendency_monitor.get_pendency_by_stage(db)
    critical = pendency_monitor.get_critical_pendencies(db, limit=50)
    combined = [{"report": "stage", **s} for s in stages] + [{"report": "critical", **c} for c in critical]
    return _to_csv(combined)


def export_performance_report(db: Session) -> str:
    return _to_csv(performance_analytics.get_eso_performance(db))


def export_inventory_report(db: Session) -> str:
    items = inventory_intelligence.get_inventory_overview(db)
    risks = inventory_intelligence.get_shortage_risk(db)
    combined = [{"report": "stock", **i} for i in items[:100]] + [{"report": "shortage", **r} for r in risks]
    return _to_csv(combined)


def export_tat_report(db: Session) -> str:
    stages = tat_analysis.get_tat_by_stage(db)
    states = tat_analysis.get_tat_by_state(db)
    return _to_csv([{"type": "stage", **s} for s in stages] + [{"type": "state", **s} for s in states])


def get_management_summary(db: Session) -> dict:
    from app.services.dashboard import get_dashboard_summary
    from app.services.ai_alerts import get_alerts

    summary = get_dashboard_summary(db)
    alerts = get_alerts(db, unresolved_only=True)[:5]
    top_states = revenue_analytics.get_revenue_by_state(db, limit=3)
    pendency = pendency_monitor.get_pendency_overview(db)

    return {
        "executive_summary": (
            f"Real Industries HSRP Operations — {summary['total_orders']} orders, "
            f"₹{summary['total_revenue']:,.0f} revenue. "
            f"{pendency['total_delayed']} orders breaching SLA. "
            f"{summary['critical_alerts']} management alerts active."
        ),
        "kpis": summary,
        "top_revenue_states": top_states,
        "priority_alerts": [
            {"title": a.title, "severity": a.severity, "message": a.message}
            for a in alerts
        ],
        "recommendations": [
            a.recommendation for a in alerts if a.recommendation
        ][:5],
    }


def export_executive_ppt(db: Session) -> bytes:
    """Generate executive summary PowerPoint for management MIS."""
    from io import BytesIO

    from pptx import Presentation
    from pptx.util import Inches, Pt

    from app.services.dashboard import get_dashboard_summary

    summary = get_dashboard_summary(db)
    mgmt = get_management_summary(db)

    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "HSRP Operations — Executive Summary"
    title.placeholders[1].text = mgmt["executive_summary"]

    kpi_slide = prs.slides.add_slide(prs.slide_layouts[1])
    kpi_slide.shapes.title.text = "Key Performance Indicators"
    body = kpi_slide.placeholders[1].text_frame
    body.clear()
    kpis = [
        ("Total Orders", summary["total_orders"]),
        ("Total Revenue", f"₹{summary['total_revenue']:,.0f}"),
        ("Pending Orders", summary["pending_orders"]),
        ("Critical Alerts", summary["critical_alerts"]),
        ("Avg TAT (hours)", summary["avg_tat_hours"]),
    ]
    for label, value in kpis:
        p = body.add_paragraph()
        p.text = f"{label}: {value}"
        p.font.size = Pt(18)

    alert_slide = prs.slides.add_slide(prs.slide_layouts[1])
    alert_slide.shapes.title.text = "Priority Alerts"
    alert_body = alert_slide.placeholders[1].text_frame
    alert_body.clear()
    for alert in mgmt.get("priority_alerts", [])[:6]:
        p = alert_body.add_paragraph()
        p.text = f"[{alert['severity'].upper()}] {alert['title']}"
        p.font.size = Pt(14)

    state_slide = prs.slides.add_slide(prs.slide_layouts[5])
    state_slide.shapes.title.text = "Top Revenue States"
    rows = [["State", "Revenue", "Orders"]]
    for s in mgmt.get("top_revenue_states", []):
        rows.append([s["name"], f"₹{s['revenue']:,.0f}", str(s.get("order_count", ""))])
    if len(rows) > 1:
        table = state_slide.shapes.add_table(
            len(rows), 3, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * len(rows))
        ).table
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                table.cell(r, c).text = val

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
